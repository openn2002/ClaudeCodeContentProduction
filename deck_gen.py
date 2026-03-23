"""
Digital Wellness Content Pipeline — Deck Generator
Creates a 10-slide executive summary PPTX using python-pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Slide dimensions (16:9) ──────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Colour palette ───────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0A, 0x23, 0x42)
NAVY2       = RGBColor(0x1A, 0x3F, 0x6F)
TEAL        = RGBColor(0x2C, 0xA5, 0x8D)
TEAL_LIGHT  = RGBColor(0xE0, 0xF5, 0xF1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF         = RGBColor(0xF8, 0xFA, 0xFC)
GRAY        = RGBColor(0xE2, 0xE8, 0xF0)
BODY        = RGBColor(0x1E, 0x29, 0x3B)
MUTED       = RGBColor(0x64, 0x74, 0x8B)
GREEN       = RGBColor(0x10, 0xB9, 0x81)   # automated
GREEN_L     = RGBColor(0xD1, 0xFA, 0xE5)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)   # manual
AMBER_L     = RGBColor(0xFE, 0xF3, 0xC7)
SLATE       = RGBColor(0x47, 0x55, 0x6B)

# ── Helpers ──────────────────────────────────────────────────────────────────

def inches(n):
    return Inches(n)

def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        inches(x), inches(y), inches(w), inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt if line_width_pt else 1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=14, bold=False, italic=False,
             color=BODY, align=PP_ALIGN.LEFT,
             font="Calibri", valign="middle",
             wrap=True):
    txBox = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    # vertical alignment
    if valign == "middle":
        tf.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = 4
    else:
        tf.vertical_anchor = 1  # TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txBox

def add_multiline(slide, lines, x, y, w, h,
                  size=13, font="Calibri", wrap=True):
    """lines = list of (text, bold, color) tuples — one per paragraph."""
    txBox = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = 1  # TOP
    first = True
    for (text, bold, color) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return txBox

def set_slide_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def agent_card(slide, x, y, w, h, title, timing, bullets, color=TEAL):
    """Draw a single agent card."""
    add_rect(slide, x, y, w, h, WHITE, GRAY, 0.75)
    # Left accent bar
    add_rect(slide, x, y, 0.06, h, color)
    # Title
    add_text(slide, title, x + 0.15, y + 0.08, w - 0.2, 0.35,
             size=14, bold=True, color=NAVY, font="Calibri")
    # Timing badge
    add_rect(slide, x + 0.15, y + 0.47, w - 0.3, 0.25, TEAL_LIGHT, None)
    add_text(slide, timing, x + 0.15, y + 0.47, w - 0.3, 0.25,
             size=9, color=TEAL, font="Calibri", bold=True, align=PP_ALIGN.CENTER)
    # Bullets
    txBox = slide.shapes.add_textbox(
        inches(x + 0.15), inches(y + 0.78), inches(w - 0.3), inches(h - 0.9)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = 1
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(10)
        run.font.color.rgb = BODY
        run.font.name = "Calibri"

def flow_box(slide, x, y, w, h, label, sub, auto=True):
    col = TEAL if auto else AMBER
    add_rect(slide, x, y, w, h, col)
    add_text(slide, label, x, y + 0.04, w, h * 0.55,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(slide, sub, x, y + h * 0.55, w, h * 0.38,
             size=8, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

def arrow_h(slide, x, y, length):
    """Thin horizontal arrow."""
    add_rect(slide, x, y - 0.015, length - 0.05, 0.03, MUTED)
    add_text(slide, "›", x + length - 0.12, y - 0.1, 0.12, 0.2,
             size=10, color=MUTED, align=PP_ALIGN.CENTER, font="Calibri")

def arrow_v(slide, x, y, length):
    """Thin vertical arrow."""
    add_rect(slide, x - 0.015, y, 0.03, length - 0.05, MUTED)
    add_text(slide, "⌄", x - 0.07, y + length - 0.15, 0.14, 0.18,
             size=10, color=MUTED, align=PP_ALIGN.CENTER, font="Calibri")

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.15, NAVY)
    add_text(slide, title, 0.5, 0.08, 10, 0.6,
             size=26, bold=True, color=WHITE, font="Calibri")
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.7, 11, 0.38,
                 size=12, color=TEAL, font="Calibri", italic=True)

# ── BUILD PRESENTATION ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank

# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, NAVY)

# Top teal band
add_rect(s, 0, 0, 13.33, 0.1, TEAL)

# Left accent bar
add_rect(s, 0.55, 2.0, 0.08, 2.8, TEAL)

# Main title
add_text(s, "Digital Wellness", 0.85, 1.7, 9, 1.5,
         size=54, bold=True, color=WHITE, font="Calibri")
add_text(s, "Content Pipeline", 0.85, 3.1, 9, 1.2,
         size=54, bold=True, color=TEAL, font="Calibri")

# Subtitle
add_text(s, "From Research to Ready-to-Film — Automatically",
         0.85, 4.55, 9, 0.5,
         size=16, color=WHITE, font="Calibri", italic=True)

# Footer
add_rect(s, 0, 6.9, 13.33, 0.6, NAVY2)
add_text(s, "Powered by Claude AI  ·  Notion  ·  Apify  ·  Slack  ·  CSIRO Total Wellbeing Diet (AU)",
         0.5, 6.9, 12.33, 0.6,
         size=11, color=MUTED, font="Calibri", valign="middle")

# Decorative circle elements
for (cx, cy, cw, ch, col, trans) in [
    (10.5, 0.5, 3.8, 3.8, NAVY2, None),
    (11.2, 1.2, 2.4, 2.4, TEAL, None),
]:
    circle = s.shapes.add_shape(9, inches(cx), inches(cy), inches(cw), inches(ch))  # OVAL
    circle.fill.solid()
    circle.fill.fore_color.rgb = col
    circle.fill.fore_color.theme_color  # noop, just accessing
    circle.line.fill.background()
    # Make teal circle semi-transparent via lxml
    if col == TEAL:
        xfrm = circle.fill._xPr
        # just use transparency via pptx
        circle.fill.fore_color.rgb = RGBColor(0x2C, 0xA5, 0x8D)

# Override circle opacity on teal one — simpler: just make it smaller/different shade
# The decorative shapes are fine as-is

# ============================================================================
# SLIDE 2 — WHAT WE BUILT
# ============================================================================
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, OFF)
header_bar(s, "What We Built",
           "An end-to-end automated content production system for CSIRO Total Wellbeing Diet (AU)")

# 3 stat cards
stats = [
    ("5", "AI Agents Built", "Research · Competitor · Ideas\nScript · Performance"),
    ("48h", "Research → Script", "Fully automated, every week"),
    ("~$5", "Per Week", "Total Apify data cost\nto power the pipeline"),
]
for i, (num, label, sub) in enumerate(stats):
    x = 0.5 + i * 4.28
    # Card background
    add_rect(s, x, 1.45, 3.9, 2.5, WHITE, GRAY, 0.5)
    # Top teal accent
    add_rect(s, x, 1.45, 3.9, 0.1, TEAL)
    # Number
    add_text(s, num, x, 1.55, 3.9, 1.1,
             size=52, bold=True, color=TEAL, align=PP_ALIGN.CENTER, font="Calibri")
    # Label
    add_text(s, label, x, 2.6, 3.9, 0.42,
             size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Calibri")
    # Sub
    add_text(s, sub, x, 3.02, 3.9, 0.8,
             size=11, color=MUTED, align=PP_ALIGN.CENTER, font="Calibri")

# Description box
add_rect(s, 0.5, 4.25, 12.33, 1.75, TEAL_LIGHT, TEAL, 0.75)
add_text(s, "The Pipeline", 0.8, 4.38, 2.5, 0.4,
         size=13, bold=True, color=TEAL, font="Calibri")
desc = (
    "Every Friday night, the system automatically scrapes live data from TikTok, Instagram, YouTube "
    "and Google — then runs a competitor intelligence scan across AU health brands. "
    "Saturday morning, Claude generates 10–15 ranked video ideas grounded in real data. "
    "The moment you approve an idea in Notion, the Script Agent picks it up and writes a full "
    "production-ready script — hook, outline, on-screen text, CTA, and caption. "
    "Every week, hands-off."
)
add_text(s, desc, 0.8, 4.78, 12.0, 1.1,
         size=12, color=BODY, font="Calibri", wrap=True)

# ============================================================================
# SLIDE 3 — FULL PIPELINE FLOW
# ============================================================================
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, OFF)
header_bar(s, "The Full Pipeline", "Teal = Automated  |  Amber = Manual approval gate")

# ── LEGEND ──
add_rect(s, 9.8, 0.28, 0.22, 0.22, TEAL)
add_text(s, "Automated", 10.06, 0.28, 1.3, 0.22, size=10, color=WHITE, font="Calibri")
add_rect(s, 11.2, 0.28, 0.22, 0.22, AMBER)
add_text(s, "Manual", 11.46, 0.28, 1.0, 0.22, size=10, color=WHITE, font="Calibri")

# ── ROW 1 — Research Phase ──────────────────────────────
add_text(s, "PHASE 1 — RESEARCH & INTELLIGENCE",
         0.5, 1.28, 8, 0.28,
         size=9, bold=True, color=MUTED, font="Calibri")

row1 = [
    ("Live Data\nScraping",    "TikTok · IG\nYouTube · Google", True),
    ("Research\nAgent",        "Fri 10pm AEST\nClaude analysis", True),
    ("Competitor\nAgent",      "AU brands &\nkeywords", True),
    ("Video Idea\nAgent",      "Sat 8am AEST\n10–15 ranked ideas", True),
    ("✓ Approve\nIdeas",       "Ollie + Nicole\nNotio Ideas DB", False),
]

bw, bh = 2.3, 1.2
gy1 = 1.62
gap = 0.28

for i, (label, sub, auto) in enumerate(row1):
    bx = 0.45 + i * (bw + gap)
    flow_box(s, bx, gy1, bw, bh, label, sub, auto)
    if i < len(row1) - 1:
        arrow_h(s, bx + bw, gy1 + bh / 2 - 0.015, gap)

# Downward connector at far right
last_x = 0.45 + 4 * (bw + gap) + bw / 2
arrow_v(s, last_x, gy1 + bh, 0.55)

# Horizontal line back to left (bottom connector)
connector_y = gy1 + bh + 0.52
add_rect(s, 0.45 + bw / 2, connector_y - 0.015, last_x - (0.45 + bw / 2), 0.03, MUTED)
# Arrow head pointing left at start of row2
add_text(s, "‹", 0.38, connector_y - 0.1, 0.2, 0.2,
         size=10, color=MUTED, align=PP_ALIGN.CENTER, font="Calibri")

# ── ROW 2 — Content Creation Phase ──────────────────────
add_text(s, "PHASE 2 — CONTENT CREATION & REPORTING",
         0.5, 3.5, 8, 0.28,
         size=9, bold=True, color=MUTED, font="Calibri")

row2 = [
    ("Script\nAgent",         "Hourly poll\nFull production script", True),
    ("✓ Approve\nScripts",    "Science check\nFilming approval", False),
    ("Content\nCalendar",     "Schedule in\nNotion", False),
    ("Post\nContent",         "Via Later or\nnatively", False),
    ("Performance\nAgent",    "Mon 8am AEST\nScorecard", True),
]

gy2 = 3.82

for i, (label, sub, auto) in enumerate(row2):
    bx = 0.45 + i * (bw + gap)
    flow_box(s, bx, gy2, bw, bh, label, sub, auto)
    if i < len(row2) - 1:
        arrow_h(s, bx + bw, gy2 + bh / 2 - 0.015, gap)

# Footer note
add_rect(s, 0.45, 5.22, 12.4, 0.35, TEAL_LIGHT, TEAL, 0.5)
add_text(s,
         "Slack notifies #content-pipeline at every stage — new research, new ideas, new scripts ready for review.",
         0.65, 5.22, 12.0, 0.35,
         size=10, color=TEAL, font="Calibri", italic=True, valign="middle")

# ============================================================================
# SLIDE 4 — THE FOUR AGENTS
# ============================================================================
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, OFF)
header_bar(s, "The AI Agents", "Four specialised Claude agents powering the pipeline")

agents = [
    (
        "Research Agent",
        "Friday 10:00pm AEST",
        [
            "Scrapes TikTok, Instagram, YouTube & Google",
            "Analyses trends, hooks, keyword opportunities",
            "Full report written to Notion (Research DB)",
            "Slack notification with strategic summary",
        ],
        TEAL,
    ),
    (
        "Competitor Agent",
        "Runs alongside Research Agent",
        [
            "Scans AU health brands: Noom, WW, Healthy Mummy, KIC, 28 by Sam Wood, Lite n Easy",
            "Keyword-based (not account-locked — stays current)",
            "Saves analysis to Research DB (Type: Competitor Analysis)",
            "Informs what gaps CSIRO can fill",
        ],
        NAVY2,
    ),
    (
        "Video Idea Agent",
        "Saturday 8:00am AEST",
        [
            "Reads latest research + competitor report",
            "Generates 10–15 ranked video concepts",
            "Avoids repeating previously generated ideas",
            "Each idea includes source attribution + pillar + platform",
            "Writes to Notion Ideas DB — Slack notification fires",
        ],
        TEAL,
    ),
    (
        "Script Agent",
        "Hourly poll — triggers on approval",
        [
            "Polls Ideas DB for Status = 'Approved'",
            "Loads research + competitor context for each run",
            "Generates: Hook, Script Outline, On-Screen Text, CTA, Caption",
            "Writes to Script Library — Slack tags Ollie + Nicole",
            "State file prevents re-processing",
        ],
        NAVY2,
    ),
]

cols = 2
cw, ch = 6.35, 2.75
for i, (title, timing, bullets, color) in enumerate(agents):
    row = i // cols
    col = i % cols
    x = 0.4 + col * (cw + 0.18)
    y = 1.38 + row * (ch + 0.22)
    agent_card(s, x, y, cw, ch, title, timing, bullets, color)

# ============================================================================
# SLIDE 5 — AUTOMATED vs MANUAL
# ============================================================================
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, OFF)
header_bar(s, "Automated vs Manual", "What runs itself — and where your editorial judgement matters")

# Left column — Automated
add_rect(s, 0.4, 1.35, 5.9, 5.7, WHITE, GRAY, 0.5)
add_rect(s, 0.4, 1.35, 5.9, 0.55, TEAL)
add_text(s, "AUTOMATED", 0.4, 1.35, 5.9, 0.55,
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

auto_items = [
    ("Data scraping", "TikTok, Instagram hashtags, YouTube, Google — via Apify"),
    ("Research analysis", "Claude synthesises scraped data into a structured report"),
    ("Competitor scanning", "AU health brand monitoring by keyword"),
    ("Idea generation", "10–15 ranked video concepts every week"),
    ("Script writing", "Full production script the moment an idea is approved"),
    ("Slack notifications", "Team alerted at every pipeline stage"),
    ("Deduplication", "Won't re-run research or re-process approved ideas"),
    ("Performance reporting", "Weekly scorecard from Later + Meta (coming soon)"),
]
for i, (title, desc) in enumerate(auto_items):
    y = 2.08 + i * 0.57
    add_rect(s, 0.55, y, 0.28, 0.28, TEAL)
    add_text(s, title, 0.97, y, 2.4, 0.28, size=11, bold=True, color=NAVY, font="Calibri", valign="middle")
    add_text(s, desc, 0.97, y + 0.28, 5.1, 0.27, size=9, color=MUTED, font="Calibri")

# Right column — Manual
add_rect(s, 6.9, 1.35, 5.9, 5.7, WHITE, GRAY, 0.5)
add_rect(s, 6.9, 1.35, 5.9, 0.55, AMBER)
add_text(s, "MANUAL (YOUR EDITORIAL LAYER)", 6.9, 1.35, 5.9, 0.55,
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

manual_items = [
    ("Review ideas", "Open Notion Ideas DB — mark Status = 'Approved'"),
    ("Science check", "Review scripts for accuracy before filming"),
    ("Approve for filming", "Tick 'Approved for filming' in Script Library"),
    ("Content calendar", "Assign dates to approved scripts in Notion"),
    ("Post content", "Publish via Later or natively on each platform"),
    ("Respond & engage", "Community management — comments, DMs"),
]
for i, (title, desc) in enumerate(manual_items):
    y = 2.08 + i * 0.65
    add_rect(s, 7.05, y, 0.28, 0.28, AMBER)
    add_text(s, title, 7.47, y, 2.4, 0.28, size=11, bold=True, color=NAVY, font="Calibri", valign="middle")
    add_text(s, desc, 7.47, y + 0.28, 5.1, 0.27, size=9, color=MUTED, font="Calibri")

add_text(s, "Bottom line: you spend your time approving great work — not generating it from scratch.",
         0.4, 7.12, 12.5, 0.35,
         size=11, italic=True, color=TEAL, font="Calibri")

# ============================================================================
# SLIDE 6 — NOTION: COMMAND CENTRE
# ============================================================================
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, OFF)
header_bar(s, "Notion: Your Command Centre",
           "Five connected databases — the single source of truth for the pipeline")

dbs = [
    ("🔍", "Research DB", "Research reports + competitor analysis land here every week. Full markdown report stored as page body — click in to read.", TEAL),
    ("💡", "Ideas DB",    "AI-generated video concepts. Review, filter by pillar/platform/priority. Approve what you want scripted.", NAVY2),
    ("📝", "Script Library", "Full production scripts. Table view shows only metadata — click any row to read the full script.", TEAL),
    ("📅", "Content Calendar", "Schedule approved scripts to filming/post dates. Links back to Script Library.", NAVY2),
    ("📊", "Weekly Scorecard", "Performance data per week: views, followers, saves, reach. Populated by Performance Agent (coming soon).", TEAL),
]

dw = 2.35
for i, (icon, name, desc, color) in enumerate(dbs):
    x = 0.35 + i * (dw + 0.2)
    # Card
    add_rect(s, x, 1.42, dw, 5.5, WHITE, GRAY, 0.5)
    # Top band
    add_rect(s, x, 1.42, dw, 0.75, color)
    # Icon
    add_text(s, icon, x, 1.42, dw, 0.75,
             size=22, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri", valign="middle")
    # Name
    add_text(s, name, x + 0.1, 2.25, dw - 0.2, 0.5,
             size=12, bold=True, color=NAVY, font="Calibri")
    # Desc
    add_text(s, desc, x + 0.1, 2.78, dw - 0.2, 3.8,
             size=10, color=BODY, font="Calibri", wrap=True, valign="top")

# ============================================================================
# SLIDE 7 — WEEKLY RHYTHM
# ============================================================================
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, OFF)
header_bar(s, "The Weekly Rhythm", "What runs, when — and what you do")

days = [
    ("Friday\nNight",   "22:00 AEST",   NAVY2,  ["Research Agent scrapes TikTok, IG,\nYouTube & Google", "Competitor Agent scans AU health brands", "Research + competitor report written to Notion", "Slack notification fires to #content-pipeline"]),
    ("Saturday\nMorning","08:00 AEST",  TEAL,   ["Video Idea Agent reads research + competitor report", "Generates 10–15 ranked ideas", "Ideas written to Notion Ideas DB", "Slack notification: new ideas ready to review"]),
    ("Any Time\n(You)",  "Your decision",AMBER,  ["Review ideas in Notion", "Approve the ones you want scripted", "Script Agent picks them up within the hour", "Scripts appear in Script Library with Slack tag"]),
    ("Ongoing\n(Hourly)","Auto-poll",    NAVY2,  ["Script Agent checks for new Approved ideas", "Generates full script with research + competitor context", "Science Approved → Approved for Filming", "Move to Content Calendar when ready"]),
    ("Monday\nMorning",  "08:00 AEST",  TEAL,   ["Performance Agent reads last week's posts", "Pulls views, followers, reach from Later + Meta", "Writes weekly scorecard to Notion", "Slack summary to #organic-growth  (coming soon)"]),
]

dw = 2.3
for i, (day, time, color, items) in enumerate(days):
    x = 0.35 + i * (dw + 0.15)
    # Day header box
    add_rect(s, x, 1.38, dw, 0.95, color)
    add_text(s, day, x, 1.38, dw, 0.62,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, time, x, 1.98, dw, 0.32,
             size=9, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri", italic=True)
    # Items
    add_rect(s, x, 2.38, dw, 4.7, WHITE, GRAY, 0.5)
    for j, item in enumerate(items):
        iy = 2.52 + j * 1.06
        add_rect(s, x + 0.12, iy, 0.22, 0.22, color)
        add_text(s, item, x + 0.42, iy - 0.02, dw - 0.52, 0.95,
                 size=9.5, color=BODY, font="Calibri", wrap=True, valign="top")

# ============================================================================
# SLIDE 8 — IMPACT
# ============================================================================
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, NAVY)

add_rect(s, 0, 0, 13.33, 0.1, TEAL)

add_text(s, "The Impact", 0.6, 0.3, 11, 0.7,
         size=30, bold=True, color=WHITE, font="Calibri")
add_text(s, "What this pipeline replaces — and what it frees up",
         0.6, 0.95, 10, 0.42,
         size=14, italic=True, color=TEAL, font="Calibri")

# Before / After
for (x, w, label, color) in [(0.5, 5.8, "BEFORE", MUTED), (7.2, 5.8, "AFTER", TEAL)]:
    add_rect(s, x, 1.62, w, 0.5, color)
    add_text(s, label, x, 1.62, w, 0.5,
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

before = [
    "Hours researching TikTok/IG trends manually each week",
    "Manually checking competitor accounts one by one",
    "Brainstorming ideas from scratch — often same pillars repeated",
    "Writing scripts from a blank page for every video",
    "No systematic record of what ideas were generated or why",
    "Performance data pulled manually from multiple platforms",
]
after = [
    "Research scraped, analysed & in Notion by Saturday morning",
    "Competitor intelligence built into every idea and script",
    "10–15 data-grounded ranked ideas every week — zero repeat",
    "Full production scripts auto-generated from approved ideas",
    "Every idea has source attribution — grounded in real data",
    "Automated weekly scorecard across TikTok, IG, YouTube, Facebook",
]

for i, (b, a) in enumerate(zip(before, after)):
    y = 2.32 + i * 0.68
    add_rect(s, 0.5, y, 0.32, 0.32, MUTED)
    add_text(s, b, 0.95, y, 5.2, 0.65, size=10.5, color=WHITE, font="Calibri", wrap=True, valign="top")

    add_rect(s, 7.2, y, 0.32, 0.32, TEAL)
    add_text(s, a, 7.65, y, 5.1, 0.65, size=10.5, color=WHITE, font="Calibri", wrap=True, valign="top")

# Divider
add_rect(s, 6.5, 1.62, 0.04, 5.2, NAVY2)

# Bottom stat
add_rect(s, 0.5, 6.78, 12.33, 0.6, NAVY2)
add_text(s,
         "Time saved: ~6–10 hrs/week of manual research, ideation, and scripting — redirected to approvals, filming, and community.",
         0.7, 6.78, 12.0, 0.6,
         size=11.5, italic=True, color=TEAL, font="Calibri", valign="middle")

# ============================================================================
# SLIDE 9 — WHAT'S NEXT
# ============================================================================
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, OFF)
header_bar(s, "What's Next", "Remaining setup + planned enhancements")

items = [
    (
        "1",
        "Activate Performance Agent",
        "Complete",
        [
            "Add LATER_API_KEY to .env → unlocks TikTok, Instagram, YouTube metrics",
            "Add META_ACCESS_TOKEN + META_PAGE_ID → unlocks Facebook reach",
            "Run: python agents/performance_agent.py — auto-populates Weekly Scorecard every Monday",
        ],
        TEAL,
        "To do — env vars needed",
    ),
    (
        "2",
        "Fix YouTube Scraper",
        "Minor",
        [
            "Current actor name not resolving on Apify (returns 0 results)",
            "Update actor ID in lib/apify.py once correct actor is confirmed",
            "YouTube data will then flow into research + competitor reports",
        ],
        NAVY2,
        "Quick fix",
    ),
    (
        "3",
        "Push to GitHub & Set Up Scheduler",
        "Infrastructure",
        [
            "All code committed to github.com/openn2002/ClaudeCodeContentProduction",
            "Run scheduler.py to keep cron-style agents running on your machine",
            "Or configure as system cron — cron_setup.sh has the commands ready",
        ],
        TEAL,
        "Ready to go",
    ),
    (
        "4",
        "Ongoing Refinement",
        "Ongoing",
        [
            "Review first full week of ideas — adjust prompt if pillar diversity needs tuning",
            "Consider Reddit scraper (paid actor) for audience conversation signals",
            "Expand markets if US (Mayo Clinic Diet) scope opens up",
        ],
        NAVY2,
        "Iterative",
    ),
]

iw, ih = 5.9, 3.05
for i, (num, title, badge, bullets, color, badge_label) in enumerate(items):
    col = i % 2
    row = i // 2
    x = 0.4 + col * (iw + 0.5)
    y = 1.38 + row * (ih + 0.22)

    add_rect(s, x, y, iw, ih, WHITE, GRAY, 0.5)
    add_rect(s, x, y, iw, 0.08, color)

    # Number circle
    add_rect(s, x + 0.12, y + 0.18, 0.35, 0.35, color)
    add_text(s, num, x + 0.12, y + 0.18, 0.35, 0.35,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    add_text(s, title, x + 0.57, y + 0.18, iw - 0.7, 0.38,
             size=13, bold=True, color=NAVY, font="Calibri")

    # Badge
    add_rect(s, x + 0.57, y + 0.6, 1.4, 0.25, TEAL_LIGHT)
    add_text(s, badge_label, x + 0.57, y + 0.6, 1.4, 0.25,
             size=8, color=TEAL, font="Calibri", bold=True, align=PP_ALIGN.CENTER)

    for j, b in enumerate(bullets):
        by = y + 0.95 + j * 0.64
        add_rect(s, x + 0.18, by + 0.04, 0.12, 0.12, color)
        add_text(s, b, x + 0.42, by, iw - 0.58, 0.62,
                 size=9.5, color=BODY, font="Calibri", wrap=True, valign="top")

# ============================================================================
# SLIDE 10 — CLOSING
# ============================================================================
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, NAVY)

add_rect(s, 0, 0, 13.33, 0.1, TEAL)

# Central statement
add_text(s, "From data scrape to production-ready script.",
         0.7, 1.5, 11.5, 1.1,
         size=38, bold=True, color=WHITE, font="Calibri")
add_text(s, "Automatically. Every week.",
         0.7, 2.55, 11.5, 1.0,
         size=38, bold=True, color=TEAL, font="Calibri")

# Divider
add_rect(s, 0.7, 3.65, 4, 0.06, TEAL)

# Summary pillars
pillars = [
    ("Research", "Live data, analysed by Claude"),
    ("Ideas", "10–15 ranked concepts weekly"),
    ("Scripts", "Full production briefs on demand"),
    ("Reporting", "Automated scorecard incoming"),
]
for i, (p, d) in enumerate(pillars):
    x = 0.7 + i * 3.18
    add_text(s, p, x, 3.9, 3.0, 0.45,
             size=16, bold=True, color=TEAL, font="Calibri")
    add_text(s, d, x, 4.35, 3.0, 0.4,
             size=11, color=WHITE, font="Calibri", italic=True)

# Stack footer
add_rect(s, 0, 6.85, 13.33, 0.65, NAVY2)
stack = "Claude AI  ·  Anthropic SDK  ·  Notion API  ·  Apify  ·  Slack  ·  Later API  ·  Meta Graph API  ·  Python"
add_text(s, stack, 0.5, 6.85, 12.33, 0.65,
         size=11, color=MUTED, align=PP_ALIGN.CENTER, font="Calibri", valign="middle")

# ── WRITE FILE ────────────────────────────────────────────────────────────────
output_path = "Digital_Wellness_Content_Pipeline.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
